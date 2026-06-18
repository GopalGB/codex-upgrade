#!/usr/bin/env python3
"""deck.py - read and generate PowerPoint decks (no MCP) via python-pptx.

(Named deck.py, NOT pptx.py, on purpose: a module named pptx.py would shadow the
installed python-pptx package and break `from pptx import Presentation`.)

Subcommands:
  dump  FILE.pptx                       extract all text, slide by slide (markdown)
  info  FILE.pptx                       slide count, layout names, image/table counts
  gen   OUT.pptx --spec SPEC.(md|json)  build a deck from a simple spec

Spec formats for `gen`:
  Markdown:  "# Title" starts a new slide (title); "- bullet" adds body bullets;
             a line starting "## " becomes the subtitle of the current slide.
  JSON:      [{"title": "...", "subtitle": "...", "bullets": ["a","b"]}, ...]

Examples:
  python deck.py dump deck.pptx
  python deck.py gen out.pptx --spec outline.md
"""

from __future__ import annotations

import argparse
import json
import os
import sys
from pathlib import Path

_HERE = os.path.dirname(os.path.abspath(__file__))
for _c in (
    os.path.join(_HERE, "..", "..", "..", "lib"),
    os.path.expanduser("~/.codex/lib"),
    os.path.expanduser("~/.agents/lib"),
):
    if os.path.isfile(os.path.join(_c, "codex_env.py")):
        sys.path.insert(0, os.path.abspath(_c))
        break
import codex_env  # noqa: E402


def _open(path: str):
    from pptx import Presentation

    return Presentation(path)


def cmd_dump(a):
    prs = _open(a.file)
    for i, slide in enumerate(prs.slides, start=1):
        print(f"\n## Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = "".join(r.text for r in para.runs) or para.text
                    if txt.strip():
                        indent = "  " * getattr(para, "level", 0)
                        print(f"{indent}- {txt.strip()}")
            if shape.has_table:
                tbl = shape.table
                for row in tbl.rows:
                    print("  | " + " | ".join(c.text for c in row.cells) + " |")


def cmd_info(a):
    prs = _open(a.file)
    print(f"file: {a.file} ({Path(a.file).stat().st_size / 1024:.0f}KB)")
    print(f"slides: {sum(1 for _ in prs.slides)}")
    print(
        f"size: {prs.slide_width} x {prs.slide_height} EMU "
        f"({prs.slide_width / 914400:.1f}in x {prs.slide_height / 914400:.1f}in)"
    )
    for i, slide in enumerate(prs.slides, start=1):
        tables = sum(1 for s in slide.shapes if s.has_table)
        pics = sum(1 for s in slide.shapes if s.shape_type == 13)
        title = ""
        for s in slide.shapes:
            if s.has_text_frame and s.text_frame.text.strip():
                title = s.text_frame.text.strip().splitlines()[0][:50]
                break
        print(f"  {i:>3}. {title:<52} tables={tables} images={pics}")


def _parse_md(text: str):
    slides, cur = [], None
    for raw in text.splitlines():
        line = raw.rstrip()
        if line.startswith("# "):
            if cur:
                slides.append(cur)
            cur = {"title": line[2:].strip(), "subtitle": "", "bullets": []}
        elif line.startswith("## ") and cur is not None:
            cur["subtitle"] = line[3:].strip()
        elif (line.lstrip().startswith(("- ", "* "))) and cur is not None:
            cur["bullets"].append(line.lstrip()[2:].strip())
        elif line.strip() and cur is not None:
            cur["bullets"].append(line.strip())
    if cur:
        slides.append(cur)
    return slides


def cmd_gen(a):
    from pptx import Presentation
    from pptx.util import Pt

    raw = Path(a.spec).read_text(encoding="utf-8")
    spec = json.loads(raw) if a.spec.endswith(".json") else _parse_md(raw)
    if not spec:
        sys.exit("spec produced 0 slides - check the format (see --help)")
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    body_layout = prs.slide_layouts[1]
    for idx, s in enumerate(spec):
        layout = title_layout if idx == 0 and not s.get("bullets") else body_layout
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = s.get("title", "")
        bullets = s.get("bullets", [])
        sub = s.get("subtitle", "")
        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                body = ph
                break
        if body is not None and (bullets or sub):
            tf = body.text_frame
            tf.clear()
            first = True
            if sub:
                tf.paragraphs[0].text = sub
                tf.paragraphs[0].font.size = Pt(20)
                first = False
            for b in bullets:
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.text = b
                p.font.size = Pt(18)
                first = False
        elif sub and slide.placeholders:
            try:
                slide.placeholders[1].text = sub
            except (KeyError, IndexError):
                pass
    prs.save(a.output)
    print(f"wrote {a.output} ({len(spec)} slides)")


def main():
    p = argparse.ArgumentParser(description="read + generate PowerPoint decks")
    sub = p.add_subparsers(dest="cmd", required=True)
    s = sub.add_parser("dump")
    s.add_argument("file")
    s.set_defaults(fn=cmd_dump)
    s = sub.add_parser("info")
    s.add_argument("file")
    s.set_defaults(fn=cmd_info)
    s = sub.add_parser("gen")
    s.add_argument("output")
    s.add_argument("--spec", required=True)
    s.set_defaults(fn=cmd_gen)
    a = p.parse_args()
    codex_env.bootstrap("python-pptx", import_names={"python-pptx": "pptx"})
    a.fn(a)


if __name__ == "__main__":
    main()
